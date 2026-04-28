import aiohttp
import aiofiles
import os
from config import FIBONACCI_API_KEY, FIBONACCI_API_URL

DESIGN_MAP = {
    "Стандартный": "standard",
    "Стильный": "stylish",
    "Креативный": "creative",
}


async def create_presentation(
    topic: str,
    language: str,
    slides_count: int,
    design: str,
    output_dir: str = "generated_files",
) -> dict[str, str]:
    """
    Вызывает Fibonacci AI API для генерации презентации.
    Возвращает словарь с путями к PDF и PPTX файлам.

    Если FIBONACCI_API_KEY не настроен — возвращает заглушку.
    """
    os.makedirs(output_dir, exist_ok=True)

    if not FIBONACCI_API_KEY:
        return await _stub_presentation(topic, output_dir)

    design_key = DESIGN_MAP.get(design, "standard")

    headers = {
        "Authorization": f"Bearer {FIBONACCI_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "topic": topic,
        "language": language,
        "slides_count": slides_count,
        "design": design_key,
    }

    async with aiohttp.ClientSession() as session:
        # Шаг 1: Создаём презентацию
        async with session.post(
            f"{FIBONACCI_API_URL}/presentations",
            headers=headers,
            json=payload,
        ) as resp:
            if resp.status not in (200, 201, 202):
                text = await resp.text()
                raise RuntimeError(f"Fibonacci API error {resp.status}: {text}")
            data = await resp.json()

        presentation_id = data.get("id") or data.get("presentation_id")
        if not presentation_id:
            raise RuntimeError("Fibonacci API: не получен ID презентации")

        # Шаг 2: Ждём и скачиваем файлы
        # NOTE: Реальный polling или webhook зависит от документации Fibonacci API.
        # Здесь простая реализация — адаптируй под реальный API.
        pdf_url = data.get("pdf_url")
        pptx_url = data.get("pptx_url")

        pdf_path = os.path.join(output_dir, f"presentation_{presentation_id}.pdf")
        pptx_path = os.path.join(output_dir, f"presentation_{presentation_id}.pptx")

        if pdf_url:
            await _download_file(session, pdf_url, pdf_path)
        if pptx_url:
            await _download_file(session, pptx_url, pptx_path)

    return {"pdf": pdf_path, "pptx": pptx_path}


async def _download_file(session: aiohttp.ClientSession, url: str, path: str):
    async with session.get(url) as resp:
        resp.raise_for_status()
        async with aiofiles.open(path, "wb") as f:
            async for chunk in resp.content.iter_chunked(8192):
                await f.write(chunk)


async def _stub_presentation(topic: str, output_dir: str) -> dict[str, str]:
    """Заглушка: создаёт пустые файлы для тестирования без реального API."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4

    safe_topic = topic[:30].replace(" ", "_").replace("/", "_")
    pdf_path = os.path.join(output_dir, f"presentation_stub_{safe_topic}.pdf")
    pptx_path = os.path.join(output_dir, f"presentation_stub_{safe_topic}.pptx")

    # Создаём заглушку PDF
    c = canvas.Canvas(pdf_path, pagesize=A4)
    c.setFont("Helvetica", 16)
    c.drawString(100, 750, f"[STUB] Презентация: {topic}")
    c.setFont("Helvetica", 12)
    c.drawString(100, 720, "Это заглушка. Настройте FIBONACCI_API_KEY.")
    c.save()

    # Создаём заглушку PPTX
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"[STUB] {topic}"
        subtitle.text = "Настройте FIBONACCI_API_KEY"
        prs.save(pptx_path)
    except ImportError:
        # python-pptx не установлен — создаём пустой файл
        async with aiofiles.open(pptx_path, "wb") as f:
            await f.write(b"STUB PPTX - install python-pptx")

    return {"pdf": pdf_path, "pptx": pptx_path}
